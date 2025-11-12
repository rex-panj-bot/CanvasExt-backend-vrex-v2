"""
File Converter - Convert Office files to PDF
Handles PPTX, DOCX, XLSX conversion to PDF format for Gemini compatibility
"""

import io
import subprocess
import tempfile
import os
from pathlib import Path
from typing import Optional, BinaryIO
import logging

logger = logging.getLogger(__name__)


def convert_office_to_pdf(file_bytes: bytes, filename: str) -> Optional[bytes]:
    """
    Convert Office/OpenDocument files to PDF

    Supported formats:
    - Microsoft Office: PPTX, DOCX, XLSX, PPT, DOC, XLS, RTF
    - OpenDocument: ODT, ODS, ODP

    Args:
        file_bytes: Original file content as bytes
        filename: Original filename (used to determine conversion method)

    Returns:
        PDF bytes if successful, None if conversion failed
    """
    ext = filename.split('.')[-1].lower() if '.' in filename else ''

    # Supported formats for PDF conversion
    supported = ['pptx', 'docx', 'xlsx', 'ppt', 'doc', 'xls', 'rtf', 'odt', 'ods', 'odp']
    if ext not in supported:
        logger.warning(f"File {filename} cannot be converted to PDF (unsupported format)")
        return None

    try:
        # Try LibreOffice conversion (most reliable)
        pdf_bytes = _convert_with_libreoffice(file_bytes, filename, ext)
        if pdf_bytes:
            return pdf_bytes

        # Fallback: Try python-pptx/docx for text extraction as PDF
        logger.info(f"LibreOffice not available, trying Python library extraction for {filename}")
        pdf_bytes = _extract_text_as_pdf(file_bytes, ext)
        return pdf_bytes

    except Exception as e:
        logger.error(f"Failed to convert {filename} to PDF: {e}")
        return None


def _convert_with_libreoffice(file_bytes: bytes, filename: str, ext: str) -> Optional[bytes]:
    """
    Convert Office file to PDF using LibreOffice headless mode
    """
    # Check if LibreOffice is available
    libreoffice_paths = [
        'libreoffice',            # Nix/Railway (in PATH)
        'soffice',                # Alternative command
        '/usr/bin/libreoffice',  # Linux
        '/usr/bin/soffice',       # Linux alternative
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # macOS
        'C:\\Program Files\\LibreOffice\\program\\soffice.exe',   # Windows
    ]

    soffice_path = None
    for path in libreoffice_paths:
        # Try direct path first
        if os.path.exists(path):
            soffice_path = path
            break
        # Try finding in PATH
        try:
            result = subprocess.run(['which', path], capture_output=True, text=True, timeout=5)
            if result.returncode == 0:
                soffice_path = path
                break
        except:
            continue

    if not soffice_path:
        print("LibreOffice not found on system")
        return None

    # Create temp directory for conversion
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, filename)

        # Write input file
        with open(input_path, 'wb') as f:
            f.write(file_bytes)

        # Convert to PDF using LibreOffice
        try:
            result = subprocess.run([
                soffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', tmpdir,
                input_path
            ], capture_output=True, text=True, timeout=30)

            if result.returncode != 0:
                logger.error(f"LibreOffice conversion failed: {result.stderr}")
                return None

            # Read converted PDF
            pdf_filename = filename.rsplit('.', 1)[0] + '.pdf'
            pdf_path = os.path.join(tmpdir, pdf_filename)

            if not os.path.exists(pdf_path):
                logger.error(f"PDF output not found: {pdf_path}")
                return None

            with open(pdf_path, 'rb') as f:
                pdf_bytes = f.read()

            logger.info(f"✅ Converted {filename} to PDF using LibreOffice ({len(pdf_bytes)} bytes)")
            return pdf_bytes

        except subprocess.TimeoutExpired:
            logger.error(f"LibreOffice conversion timeout for {filename}")
            return None
        except Exception as e:
            logger.error(f"LibreOffice conversion error: {e}")
            return None


def _extract_text_as_pdf(file_bytes: bytes, ext: str) -> Optional[bytes]:
    """
    Extract text from Office file and create simple PDF
    Fallback method when LibreOffice is not available
    """
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch

        text_content = None

        # Extract text based on file type
        if ext in ['pptx', 'ppt']:
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                slides_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text = f"Slide {i+1}:\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text += shape.text + "\n"
                    slides_text.append(slide_text)
                text_content = "\n\n".join(slides_text)
            except ImportError:
                logger.error("python-pptx not installed")
                return None

        elif ext in ['docx', 'doc']:
            try:
                from docx import Document
                doc = Document(io.BytesIO(file_bytes))
                paragraphs = [para.text for para in doc.paragraphs]
                text_content = "\n\n".join(paragraphs)
            except ImportError:
                logger.error("python-docx not installed")
                return None

        elif ext in ['xlsx', 'xls']:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes))
                sheets_text = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    sheet_text = f"Sheet: {sheet_name}\n"
                    for row in sheet.iter_rows(values_only=True):
                        row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        sheet_text += row_text + "\n"
                    sheets_text.append(sheet_text)
                text_content = "\n\n".join(sheets_text)
            except ImportError:
                logger.error("openpyxl not installed")
                return None

        if not text_content:
            return None

        # Create PDF from text
        pdf_buffer = io.BytesIO()
        doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Split text into paragraphs and add to PDF
        for para_text in text_content.split('\n\n'):
            if para_text.strip():
                para = Paragraph(para_text.replace('\n', '<br/>'), styles['Normal'])
                story.append(para)
                story.append(Spacer(1, 0.2*inch))

        doc.build(story)
        pdf_bytes = pdf_buffer.getvalue()

        logger.info(f"✅ Extracted text and created PDF ({len(pdf_bytes)} bytes)")
        return pdf_bytes

    except Exception as e:
        logger.error(f"Text extraction failed: {e}")
        return None


def convert_to_text(file_bytes: bytes, filename: str) -> Optional[bytes]:
    """
    Convert web/data formats to plain text for Gemini compatibility

    Handles: HTML, XML, JSON
    Returns plain text bytes that can be uploaded as text/plain to Gemini

    Args:
        file_bytes: Original file content as bytes
        filename: Original filename (used to determine format)

    Returns:
        Plain text bytes if successful, None if conversion failed
    """
    ext = filename.split('.')[-1].lower() if '.' in filename else ''

    try:
        # Decode bytes to string
        try:
            text_content = file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            # Try other encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    text_content = file_bytes.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                logger.error(f"Could not decode {filename} - unsupported encoding")
                return None

        if ext in ['html', 'htm']:
            # Strip HTML tags for cleaner text
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(text_content, 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                clean_text = soup.get_text(separator='\n', strip=True)
            except ImportError:
                # Fallback: basic tag stripping with regex
                import re
                clean_text = re.sub(r'<[^>]+>', '', text_content)
                clean_text = re.sub(r'\s+', ' ', clean_text).strip()

            text_content = clean_text

        elif ext == 'json':
            # Pretty-print JSON for readability
            import json
            try:
                json_obj = json.loads(text_content)
                text_content = json.dumps(json_obj, indent=2, ensure_ascii=False)
            except json.JSONDecodeError:
                # Keep original if invalid JSON
                pass

        # XML is already text, just clean up whitespace
        elif ext == 'xml':
            import re
            text_content = re.sub(r'>\s+<', '><', text_content)

        # Convert back to bytes
        result_bytes = text_content.encode('utf-8')
        logger.info(f"✅ Converted {filename} to text ({len(result_bytes)} bytes)")
        return result_bytes

    except Exception as e:
        logger.error(f"Failed to convert {filename} to text: {e}")
        return None


def needs_conversion(filename: str) -> bool:
    """Check if file needs conversion before Gemini upload"""
    ext = filename.split('.')[-1].lower() if '.' in filename else ''
    # Expanded to include OpenDocument and web formats
    return ext in [
        'pptx', 'docx', 'xlsx', 'ppt', 'doc', 'xls', 'rtf',  # MS Office
        'odt', 'ods', 'odp',  # OpenDocument
        'html', 'htm', 'xml', 'json'  # Web/data formats
    ]
